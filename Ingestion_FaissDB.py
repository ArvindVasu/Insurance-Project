import os
import json
import shutil # For potential future use, though Unstructured handles extraction

from langchain.document_loaders import (
    UnstructuredWordDocumentLoader,
    TextLoader,
    UnstructuredPowerPointLoader
)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain.embeddings.openai import OpenAIEmbeddings
from dotenv import load_dotenv
from langchain.docstore.document import Document # Ensure this is imported
from docx import Document as DocxDocument
from pptx import Presentation
from langchain.docstore.document import Document


load_dotenv()

# === CONFIG ===
DOC_FOLDER = "/Users/hp/OneDrive/Desktop/Python/Agentic AI - Vanna_Serp_Doc/Documents/"
FAISS_DIR = "faiss_index/"
EXTRACTED_IMAGES_DIR = "extracted_images/" # New folder to store images extracted from documents
IMAGE_METADATA_FILE = os.path.join(EXTRACTED_IMAGES_DIR, "extracted_image_metadata.json")

embedding_function = OpenAIEmbeddings()

# LOADER_MAP now configures for text and potentially image extraction via Unstructured
LOADER_MAP = {
    ".docx": UnstructuredWordDocumentLoader,
    ".txt": TextLoader,
    ".pptx": UnstructuredPowerPointLoader,
}


def extract_images_fallback_docx(file_path, output_dir, original_filename):
    rels = DocxDocument(file_path).part._rels
    fallback_images = []

    for i, rel in enumerate(rels.values()):
        if "image" in rel.target_ref:
            image_data = rel.target_part.blob
            ext = rel.target_part.content_type.split("/")[-1]  # e.g., png
            filename = f"{original_filename}_fallback_docx_img_{i}.{ext}"
            stored_path = os.path.join(output_dir, filename)

            with open(stored_path, "wb") as f:
                f.write(image_data)

            fallback_images.append({
                "filename": filename,
                "original_doc": original_filename,
                "original_doc_path": os.path.join("Documents", original_filename).replace("\\", "/"),
                "extracted_image_path": stored_path,
                "caption": "",  # No caption from fallback
                "image_id": f"fallback_docx_{i}",
                "page_number": None,
                "doc_type": "extracted_image"
            })

    return fallback_images


def extract_images_fallback_pptx(file_path, output_dir, original_filename):
    prs = Presentation(file_path)
    fallback_images = []
    count = 0

    for slide_index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                image = shape.image
                ext = image.ext
                image_bytes = image.blob
                filename = f"{original_filename}_fallback_pptx_img_{count}.{ext}"
                stored_path = os.path.join(output_dir, filename)

                with open(stored_path, "wb") as f:
                    f.write(image_bytes)

                fallback_images.append({
                    "filename": filename,
                    "original_doc": original_filename,
                    "original_doc_path": os.path.join("Documents", original_filename).replace("\\", "/"),
                    "extracted_image_path": stored_path,
                    "caption": "",
                    "image_id": f"fallback_pptx_{count}",
                    "page_number": slide_index + 1,
                    "doc_type": "extracted_image"
                })
                count += 1

    return fallback_images


def load_and_process_documents_with_images(folder_path):
    all_text_documents = []
    all_extracted_image_metadata = []

    os.makedirs(EXTRACTED_IMAGES_DIR, exist_ok=True)

    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        ext = os.path.splitext(filename)[1].lower()

        if ext not in LOADER_MAP:
            print(f"⏭️ Skipping unsupported file: {filename}")
            continue

        loader_class = LOADER_MAP[ext]

        try:
            if ext in [".docx", ".pptx"]:
                loader = loader_class(file_path, strategy="hi_res",
                                      extract_images_path=EXTRACTED_IMAGES_DIR)
            else:
                loader = loader_class(file_path)

            elements = loader.load()

            # Text processing
            for doc_element in elements:
                if doc_element.page_content.strip():
                    doc_element.metadata.update({
                        "source_doc": filename,
                        "extension": ext,
                        "file_path": os.path.join("Documents", filename).replace("\\", "/"),
                        "doc_type": "text_chunk"
                    })
                    all_text_documents.append(doc_element)

            # Image metadata (via Unstructured)
            image_count = 0
            for element in elements:
                if element.metadata.get("filetype", "").startswith("image/"):
                    image_filename = os.path.basename(element.metadata["image_path"])
                    stored_image_path = os.path.join(EXTRACTED_IMAGES_DIR, image_filename)

                    image_meta = {
                        "filename": image_filename,
                        "original_doc": filename,
                        "original_doc_path": os.path.join("Documents", filename).replace("\\", "/"),
                        "extracted_image_path": stored_image_path,
                        "caption": element.metadata.get("caption", ""),
                        "image_id": element.metadata.get("image_id", ""),
                        "page_number": element.metadata.get("page_number", None),
                        "doc_type": "extracted_image"
                    }
                    all_extracted_image_metadata.append(image_meta)
                    image_count += 1

            # Fallback if Unstructured missed images
            if image_count == 0:
                if ext == ".docx":
                    fallback = extract_images_fallback_docx(file_path, EXTRACTED_IMAGES_DIR, filename)
                    all_extracted_image_metadata.extend(fallback)
                    if fallback:
                        print(f"🛠️ Fallback: Extracted {len(fallback)} images from {filename}")
                elif ext == ".pptx":
                    fallback = extract_images_fallback_pptx(file_path, EXTRACTED_IMAGES_DIR, filename)
                    all_extracted_image_metadata.extend(fallback)
                    if fallback:
                        print(f"🛠️ Fallback: Extracted {len(fallback)} images from {filename}")

            print(f"✅ Processed: {filename} (Text chunks: {len(all_text_documents)}, Images: {len(all_extracted_image_metadata)})")

        except Exception as e:
            print(f"❌ Error processing {filename}: {e}")

    with open(IMAGE_METADATA_FILE, 'w') as f:
        json.dump(all_extracted_image_metadata, f, indent=4)
    print(f"✅ Saved extracted image metadata to: '{IMAGE_METADATA_FILE}'")

    return all_text_documents, all_extracted_image_metadata

def ingest_to_faiss(docs):
    if not docs:
        print("⚠️ No valid text documents found for FAISS ingestion.")
        return

    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
    chunks = []

    for original_doc in docs:
        split_chunks = splitter.split_text(original_doc.page_content)
        for chunk_text in split_chunks:
            chunk = Document(
                page_content=chunk_text,
                metadata={
                    "source_doc": original_doc.metadata.get("source_doc"),
                    "extension": original_doc.metadata.get("extension"),
                    "file_path": str(original_doc.metadata.get("file_path") or original_doc.metadata.get("source")),
                    "doc_type": original_doc.metadata.get("doc_type", "text_chunk")
                }
            )
            chunks.append(chunk)

    print(f"🧩 Total text chunks created for FAISS: {len(chunks)}")

    vectorstore = FAISS.from_documents(
        documents=chunks,
        embedding=embedding_function
    )
    vectorstore.save_local(FAISS_DIR)
    print(f"✅ Successfully ingested text content into FAISS → Directory: '{FAISS_DIR}'")

if __name__ == "__main__":
    text_docs_for_faiss, extracted_image_metadata = load_and_process_documents_with_images(DOC_FOLDER)
    ingest_to_faiss(text_docs_for_faiss)
    print("\nIngestion process complete.")

    # You can now access extracted_image_metadata for displaying image links in your Streamlit app
    # Example of how you might load and use image_metadata later:
    # with open(IMAGE_METADATA_FILE, 'r') as f:
    #     loaded_image_metadata = json.load(f)
    # print("\nLoaded Extracted Image Metadata:")
    # for img_meta in loaded_image_metadata:
    #     print(f"  - {img_meta['filename']} from '{img_meta['original_doc']}', Stored at {img_meta['extracted_image_path']}, Caption: '{img_meta['caption']}'")