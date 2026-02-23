import streamlit as st

import os
from dotenv import load_dotenv
import re
import pandas as pd

import matplotlib.pyplot as plt
import networkx as nx
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.text.paragraph import Paragraph
from docx.table import Table
from docx.table import Table as DocxTable
import tempfile
import uuid
from docx.table import Table as DocxTable
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
import numpy as np
from pathlib import Path
from io import BytesIO
import base64
from time import sleep
from urllib.parse import urlparse
import io
import re
from difflib import get_close_matches
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import math
import logging

load_dotenv()

def _rows_cols_from_serialized(df_like):
    """
    Accepts:
      - pandas.DataFrame
      - dict with {"columns": [...], "rows": [...]}
      - list[dict] (rows only)
    Returns: (columns:list[str], rows:list[list[str]])
    """
    if df_like is None:
        return [], []
    # DataFrame
    if isinstance(df_like, pd.DataFrame):
        cols = list(df_like.columns)
        rows = df_like.to_dict(orient="records")
        return cols, [[str(row.get(c, "")) for c in cols] for row in rows]
    # {"columns": [...], "rows": [...]}
    if isinstance(df_like, dict) and "rows" in df_like:
        cols = df_like.get("columns") or []
        rows_data = df_like["rows"]
        # if columns missing, infer from first row
        if not cols and isinstance(rows_data, list) and rows_data:
            cols = list(rows_data[0].keys())
        rows = []
        for r in rows_data or []:
            if isinstance(r, dict):
                rows.append([str(r.get(c, "")) for c in cols])
            else:
                # row already list-like
                rows.append([str(v) for v in (r or [])])
        return cols, rows
    # list-of-dicts
    if isinstance(df_like, list) and (not df_like or isinstance(df_like[0], dict)):
        cols = list(df_like[0].keys()) if df_like else []
        rows = [[str(r.get(c, "")) for c in cols] for r in df_like]
        return cols, rows
    # Fallback: treat as string
    return ["value"], [[str(df_like)]]


#++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++
# ---- Visual style ----
FONT_TITLE = "Segoe UI"
FONT_BODY  = "Segoe UI"
COLOR_TITLE = RGBColor(30, 30, 30)
COLOR_TEXT  = RGBColor(45, 45, 45)
COLOR_TBL_HEADER_BG = RGBColor(230, 234, 239)
COLOR_TBL_HEADER_TX = RGBColor(20, 20, 20)
COLOR_TBL_ROW_ALT   = RGBColor(247, 249, 251)
# Subtle background and title frame
COLOR_BG            = RGBColor(248, 249, 251)   # off-white/gray
COLOR_TITLE_FRAME   = RGBColor(200, 205, 210)   # soft gray line
TITLE_FRAME_LINE_W  = Pt(2)
TITLE_FRAME_PAD_X   = Pt(6)
TITLE_FRAME_PAD_Y   = Pt(4)

# Content geometry (keeps things consistent)
SLIDE_MARGIN_LEFT   = Inches(0.6)
SLIDE_MARGIN_RIGHT  = Inches(0.6)
CONTENT_TOP         = Inches(1.7)
CONTENT_W           = Inches(10) - SLIDE_MARGIN_LEFT - SLIDE_MARGIN_RIGHT
CONTENT_H           = Inches(5.2)

# limits
TOP_N = 5
MAX_ROWS_PER_SLIDE = 12      # including header row
# MAX_COLS_PER_SLIDE = 6       # show at most 6 columns per slide; rest go to next slide

def _cell_text(cell, text, *, size=11, bold=False, align=PP_ALIGN.LEFT, color=COLOR_TEXT):
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "" if text is None else str(text)
    r.font.name = FONT_BODY
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    p.alignment = align
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    # inner padding (some python-pptx versions support these)
    for attr, val in (("margin_left", Pt(2)), ("margin_right", Pt(2)), ("margin_top", Pt(2)), ("margin_bottom", Pt(2))):
        try:
            setattr(tf, attr, val)
        except Exception:
            pass

def _fill_cell(cell, rgb):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb

def _autofit_widths(table):
    """Proportionally set column widths by character length."""
    cols = table.columns
    rows = table.rows
    weights = []
    for j in range(len(cols)):
        mx = 1
        for i in range(len(rows)):
            try:
                t = rows[i].cells[j].text or ""
            except Exception:
                t = ""
            mx = max(mx, len(t))
        weights.append(mx)
    total = float(sum(weights)) or 1.0
    for j, w in enumerate(weights):
        cols[j].width = int(CONTENT_W * (w / total))        

def _chunk(lst, n):
    for i in range(0, len(lst), n):
        yield i, lst[i:i+n]

def _apply_professional_background(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def _add_slide(prs, layout):
    slide = prs.slides.add_slide(layout)
    _apply_professional_background(slide)
    return slide

def _ensure_title_shape(slide):
    """Return a usable title shape; create one if layout lacks a title."""
    shape = slide.shapes.title
    if shape is None:
        shape = slide.shapes.add_textbox(SLIDE_MARGIN_LEFT, Inches(0.6), CONTENT_W, Inches(0.95))
    return shape

def _apply_title_text(shape, text, size=34):
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT_TITLE
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

def _add_title_frame(slide, title_shape):
    """Draw a rounded-rectangle outline around the title."""
    if title_shape is None:
        return
    left  = max(0, title_shape.left  - TITLE_FRAME_PAD_X)
    top   = max(0, title_shape.top   - TITLE_FRAME_PAD_Y)
    width = title_shape.width  + 2 * TITLE_FRAME_PAD_X
    height= title_shape.height + 2 * TITLE_FRAME_PAD_Y

    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.background()  # transparent fill
    rect.line.color.rgb = COLOR_TITLE_FRAME
    rect.line.width = TITLE_FRAME_LINE_W

def _set_title(slide, text, size=34, frame=True):
    """Use this for ALL slide titles."""
    t = _ensure_title_shape(slide)
    _apply_title_text(t, text, size=size)
    if frame:
        _add_title_frame(slide, t)

def _add_table_slide(prs, title, columns, rows, max_rows=6):
    """
    Adds a slide with a well-formatted table.
    - Shows ONLY the first TOP_N rows.
    - Never splits columns; all columns appear on the same slide.
    - Wraps text, auto-fits column widths, and adjusts font size if many columns.
    """
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    try:
        # If you have _set_title for framed titles, use it:
        _set_title(slide, title)
    except Exception:
        slide.shapes.title.text = title  # fallback if you didn't add _set_title

    # Normalize rows -> list[list[str]]
    norm_rows = []
    for r in rows or []:
        if isinstance(r, dict):
            if columns:
                norm_rows.append([("" if r.get(c) is None else str(r.get(c))) for c in columns])
            else:
                norm_rows.append([("" if v is None else str(v)) for v in r.values()])
        elif isinstance(r, (list, tuple)):
            norm_rows.append([("" if v is None else str(v)) for v in r])
        else:
            norm_rows.append([("" if r is None else str(r))])

    # Use provided columns or infer max length from data
    if columns and len(columns) > 0:
        hdr = [str(c) for c in columns]
        n_cols = len(hdr)
    else:
        n_cols = max((len(r) for r in norm_rows), default=1)
        hdr = [f"Col {i+1}" for i in range(n_cols)]

    # ---- Limit to TOP_N rows (hard cap) ----
    norm_rows = norm_rows[:TOP_N]

    # Compute how many rows we can place (header + body) per slide
    rows_per_page = max(1, MAX_ROWS_PER_SLIDE - 1)  # body rows; keep 1 for header
    total_pages = max(1, math.ceil(len(norm_rows) / rows_per_page))

    left, top, width, height = SLIDE_MARGIN_LEFT, CONTENT_TOP, CONTENT_W, Inches(3.8)

    for page_idx in range(total_pages):
        start = page_idx * rows_per_page
        chunk = norm_rows[start:start + rows_per_page]

        # Add a page badge if we actually paginated by rows (rare if TOP_N=5)
        if total_pages > 1 and page_idx > 0:
            slide = prs.slides.add_slide(layout)
            try:
                _set_title(slide, f"{title} (Page {page_idx+1}/{total_pages})")
            except Exception:
                slide.shapes.title.text = f"{title} (Page {page_idx+1}/{total_pages})"

        # Build table
        n_rows = len(chunk) + 1  # header + data
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table

        # Dynamic font sizes for many columns
        header_size = 12 if n_cols <= 7 else (11 if n_cols <= 10 else 10)
        body_size   = 11 if n_cols <= 7 else (10 if n_cols <= 10 else 9)

        # Header
        for j, h in enumerate(hdr):
            cell = table.cell(0, j)
            _cell_text(cell, h, size=header_size, bold=True, align=PP_ALIGN.CENTER, color=COLOR_TBL_HEADER_TX)
            _fill_cell(cell, COLOR_TBL_HEADER_BG)

        # Body
        for i, r in enumerate(chunk, start=1):
            for j in range(n_cols):
                val = r[j] if j < len(r) else ""
                cell = table.cell(i, j)
                _cell_text(cell, val, size=body_size, bold=False, align=PP_ALIGN.LEFT, color=COLOR_TEXT)
            # zebra striping
            if i % 2 == 1:
                for j in range(n_cols):
                    _fill_cell(table.cell(i, j), COLOR_TBL_ROW_ALT)

        # Columns: proportional widths to fit CONTENT_W
        _autofit_widths(table)

def _add_text_block(slide, top=CONTENT_TOP, height=CONTENT_H):
    """Creates a word-wrapped body text box with nice defaults."""
    box = slide.shapes.add_textbox(SLIDE_MARGIN_LEFT, top, CONTENT_W, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left   = Pt(2)
    tf.margin_right  = Pt(2)
    tf.margin_top    = Pt(2)
    tf.margin_bottom = Pt(2)
    return tf

def _add_body_paragraph(tf, text, size=14, before=4, after=6, bold=False):
    p = tf.add_paragraph() if len(tf.paragraphs) else tf.paragraphs[0]
    p.text = ""
    r = p.add_run()
    r.text = text
    r.font.name = FONT_TITLE
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(before)
    p.space_after  = Pt(after)
    p.alignment = PP_ALIGN.LEFT
    tf.word_wrap = True
    return p

#Exporting data to Powerpoint
def generate_ppt(entry) -> BytesIO:
    """
    Generate a PowerPoint for a session entry which contains `messages`:
    entry["messages"] = [{"role":"turn","user_prompt":..., "assistant_run": {...}, "timestamp":...}, ...]
    Returns BytesIO.
    """
    prs = Presentation()
    layout = prs.slide_layouts[5]

    slide = _add_slide(prs, layout)
    _set_title(slide, "Agentic AI Report", size=40)  # NEW: style + border
    tf = _add_text_block(slide)
    
    session_title = entry.get("title") or (entry.get("prompt") or "")
    _add_body_paragraph(tf, session_title, size=18, before=2, after=4)
    # slide.placeholders[1].text = f"Session: {session_title}"
    created = entry.get("created_at") or entry.get("timestamp") or ""
    if created:
        # add a small subtitle for created time if available
        try:
            # subtitle = slide.placeholders[1]
            # subtitle.text += f"\nCreated: {created}"
            _add_body_paragraph(tf, f"\n⏱ Created: {created}", size=12, before=2, after=4)
        except Exception:
            pass
    
    # If messages absent (defensive) - fallback to single-run fields (but user said messages always present)
    messages = entry.get("messages", [])
    if not messages:
        # create a synthetic single-turn message using top-level entry fields
        messages = [{
            "role": "turn",
            "user_prompt": entry.get("prompt", ""),
            "assistant_run": {
                "prompt": entry.get("prompt"),
                "route": entry.get("route"),
                "result": entry.get("result"),
                "sql_query": entry.get("sql_query"),
                "web_links": entry.get("web_links"),
                "general_summary": entry.get("general_summary"),
                "comparison_summary": entry.get("comparison_summary"),
                "chart_info": entry.get("chart_info"),
                "faiss_summary": entry.get("faiss_summary"),
                "faiss_sources": entry.get("faiss_sources"),
                "faiss_images": entry.get("faiss_images"),
                "intranet_summary": entry.get("intranet_summary"),
                "intranet_sources": entry.get("intranet_sources"),
                "intranet_doc_links": entry.get("intranet_doc_links"),
                "intranet_doc_count": entry.get("intranet_doc_count")
                },
            "timestamp": entry.get("timestamp")
        }]

    # Iterate through turns in stored order (do not change order)
    for idx, turn in enumerate(messages, start=1):
        user_prompt = turn.get("user_prompt") or ""
        timestamp = turn.get("timestamp") or ""
        assistant_run = turn.get("assistant_run") or {}

        # 1) Slide for the user prompt
        slide = _add_slide(prs, layout)
        _set_title(slide, f"Turn {idx}: User Prompt")
        tf = _add_text_block(slide)
        _add_body_paragraph(tf, user_prompt, size=14, before=2, after=4)
        _add_body_paragraph(tf, f"⏱ {timestamp}", size=10, before=2, after=4)

        # If assistant_run is empty, skip assistant slides
        if not assistant_run:
            continue

        route = assistant_run.get("route")

        # --- Document related (document route) ---
        if route == "document":
            st.write('document node')

        # --- Comparison / General summaries ---
        if assistant_run.get("comparison_summary"):
            slide = _add_slide(prs, layout)
            _set_title(slide, f"Turn {idx}: Comparison Summary")
            tf = _add_text_block(slide)

            for para in str(assistant_run.get("comparison_summary")).split("\n"):
                if para.strip():
                    _add_body_paragraph(tf, para.strip(), size=12, before=2, after=4)

        # --- SQL query slide ---
        if assistant_run.get("sql_query"):
            slide = _add_slide(prs, layout)
            _set_title(slide, f"Turn {idx}: SQL Query")
            tf = _add_text_block(slide)
            _add_body_paragraph(tf, assistant_run.get("sql_query"), size=14, before=2, after=4)            

        # --- SQL Result (if any) ---
        result = assistant_run.get("result")
        df_result = None
        if isinstance(result, list):
            try:
                df_result = pd.DataFrame(result)
            except Exception:
                df_result = None
        elif isinstance(result, pd.DataFrame):
            df_result = result

        if df_result is not None and not df_result.empty and route in ["sql", "document", "comp"]:
            # Add a table slide for SQL results, cap rows
            cols, rows = _rows_cols_from_serialized(df_result)
            if rows:
                _add_table_slide(prs, f"Turn {idx}: SQL Results", cols, rows, max_rows=6)


        if assistant_run.get("general_summary"):
            slide = _add_slide(prs, layout)
            _set_title(slide, f"Turn {idx}: General Summary")
            tf = _add_text_block(slide)
            for para in str(assistant_run.get("general_summary")).split("\n"):
                if para.strip():
                    _add_body_paragraph(tf, para.strip(), size=12, before=2, after=4)


        # --- Web links (search/comp) ---
        web_links = assistant_run.get("web_links") or assistant_run.get("result") if route == "search" else assistant_run.get("web_links")
        if web_links:
            slide = _add_slide(prs, layout)
            _set_title(slide, f"Turn {idx}: Top Web Links")
            tf = _add_text_block(slide)

            for i, item in enumerate(web_links, 1):
                # item could be tuple (markdown_link, summary) or simple string
                link_md, summary = (item[0], item[1]) if (isinstance(item, (list, tuple)) and len(item) >= 2) else (str(item), "")
                match = re.match(r"\[(.*?)\]\((.*?)\)", str(link_md))
                if match:
                    title, url = match.groups()
                else:
                    title, url = f"Link {i}", str(link_md)

                p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"{i}. {title}"
                try:
                    run.font.size = Pt(12)
                    run.hyperlink.address = url
                except Exception:
                    pass
                if summary:
                    # _add_body_paragraph(tf, f"    ↳ {str(summary)[:300]}", size=12, before=2, after=4)
                    s_p = tf.add_paragraph()
                    s_p.text = f"    ↳ {str(summary)[:300]}"
                    s_p.font.size = Pt(11)

        # --- FAISS route slides if present (assistant_run or entry-level) ---
        faiss_summary = assistant_run.get("faiss_summary")
        faiss_sources = assistant_run.get("faiss_sources") or assistant_run.get("faiss_sources", [])
        faiss_images = assistant_run.get("faiss_images") or assistant_run.get("faiss_images", [])

        if faiss_summary:
            slide = _add_slide(prs, layout)
            _set_title(slide, f"Turn {idx}: FAISS Summary")
            tf = _add_text_block(slide)
            for para in str(faiss_summary).split("\n"):
                if para.strip():
                    _add_body_paragraph(tf, para.strip(), size=12, before=2, after=4)


        if faiss_sources:
            for i, src in enumerate(faiss_sources, 1):
                try:
                    docname, snippet, path = src[0], src[1], src[2] if len(src) >= 3 else (src[0], src[1], None)
                except Exception:
                    docname, snippet, path = str(src), "", None
                slide = _add_slide(prs, layout)
                _set_title(slide, f"Turn {idx}: FAISS Source {i} - {os.path.basename(path) if path else docname}")
                tf = _add_text_block(slide)
                for para in str(snippet).split("\n"):
                    if para.strip():
                        _add_body_paragraph(tf, para.strip(), size=12, before=2, after=4)

        if faiss_images and faiss_sources:
            # Only include images from the most-similar doc (first in faiss_sources)
            top_docname = faiss_sources[0][0] if isinstance(faiss_sources[0], (list, tuple)) else faiss_sources[0]
            top_images = [img for img in faiss_images if img.get("original_doc") == top_docname]
            if top_images:
                slide = _add_slide(prs, prs.slide_layouts[5])
                _set_title(slide, f"Turn {idx}: Images from {top_docname}")
                left = Inches(0.8)
                top = Inches(2.2)
                image_width = Inches(5.5)
                spacing = Inches(0.5)
                for im_meta in top_images:
                    img_path = im_meta.get("extracted_image_path")
                    if img_path and os.path.exists(img_path):
                        slide.shapes.add_picture(img_path, left, top, width=image_width)
                        top += Inches(3.2)
                        if top > Inches(6.5):
                            top = Inches(2.2)
                            left += image_width + spacing

                    # --- INTRANET ROUTE SLIDES ---
        intranet_summary = assistant_run.get("intranet_summary")
        intranet_sources = assistant_run.get("intranet_sources") or []
        intranet_doc_links = assistant_run.get("intranet_doc_links") or []
        intranet_doc_count = assistant_run.get("intranet_doc_count")
        intranet_lob = assistant_run.get("intranet_lob")

        # 🔹 Intranet Summary Slide
        if intranet_summary:
            slide = _add_slide(prs, layout)
            _set_title(
                slide,
                f"Turn {idx}: Intranet Policy Analysis" +
                (f" ({intranet_lob})" if intranet_lob else "")
            )
            tf = _add_text_block(slide)

            if intranet_doc_count is not None:
                _add_body_paragraph(
                    tf,
                    f"Documents Analysed: {intranet_doc_count}",
                    size=11,
                    before=2,
                    after=4
                )

            for para in str(intranet_summary).split("\n"):
                if para.strip():
                    _add_body_paragraph(
                        tf,
                        para.strip(),
                        size=12,
                        before=2,
                        after=4
                    )

        #  Intranet Sources Slide
        if intranet_sources:
            slide = _add_slide(prs, layout)
            _set_title(slide, f"Turn {idx}: Intranet Sources")
            tf = _add_text_block(slide)

            for i, src in enumerate(intranet_sources, 1):
                try:
                    docname, link = src
                except Exception:
                    docname, link = str(src), None

                p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"{i}. {docname}"

                try:
                    run.font.size = Pt(12)
                    if link:
                        run.hyperlink.address = link
                except Exception:
                    pass

        # 🔹 Intranet Links Slide (Optional – if separate)
        if intranet_doc_links:
            slide = _add_slide(prs, layout)
            _set_title(slide, f"Turn {idx}: Intranet Document Links")
            tf = _add_text_block(slide)

            for i, link in enumerate(intranet_doc_links, 1):
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"Document Link {i}"

                try:
                    run.font.size = Pt(12)
                    run.hyperlink.address = link
                except Exception:
                    pass

        # --- Charts: if there is chart_info (you can expand how to render charts later) ---
        chart_info = assistant_run.get("chart_info")
        if chart_info:
            slide = _add_slide(prs, layout)
            _set_title(slide, f"Turn {idx}: Chart Info")
            _add_body_paragraph(tf, str(chart_info)[:1500], size=12, before=2, after=4)
            tf = _add_text_block(slide)

    # End: return PPT as BytesIO
    ppt_bytes = BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes